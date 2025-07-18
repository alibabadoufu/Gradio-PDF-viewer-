import os
import io
import base64
from typing import List, Tuple, Optional
from PIL import Image, ImageDraw, ImageFont
import pdf2image
from docx import Document
from pptx import Presentation
import openpyxl
from openpyxl.drawing.image import Image as ExcelImage
import tempfile

class DocumentPreviewer:
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx']
    
    def is_supported(self, file_path: str) -> bool:
        """Check if the file format is supported."""
        _, ext = os.path.splitext(file_path.lower())
        return ext in self.supported_formats
    
    def get_page_count(self, file_path: str) -> int:
        """Get the total number of pages/slides/sheets in the document."""
        _, ext = os.path.splitext(file_path.lower())
        
        try:
            if ext == '.pdf':
                return self._get_pdf_page_count(file_path)
            elif ext == '.docx':
                return self._get_docx_page_count(file_path)
            elif ext == '.pptx':
                return self._get_pptx_slide_count(file_path)
            elif ext == '.xlsx':
                return self._get_excel_sheet_count(file_path)
            else:
                return 0
        except Exception as e:
            print(f"Error getting page count for {file_path}: {e}")
            return 0
    
    def preview_page(self, file_path: str, page_number: int) -> Optional[Image.Image]:
        """Generate a preview image for a specific page/slide/sheet."""
        _, ext = os.path.splitext(file_path.lower())
        
        try:
            if ext == '.pdf':
                return self._preview_pdf_page(file_path, page_number)
            elif ext == '.docx':
                return self._preview_docx_page(file_path, page_number)
            elif ext == '.pptx':
                return self._preview_pptx_slide(file_path, page_number)
            elif ext == '.xlsx':
                return self._preview_excel_sheet(file_path, page_number)
            else:
                return None
        except Exception as e:
            print(f"Error previewing page {page_number} of {file_path}: {e}")
            return self._create_error_image(f"Error loading page {page_number}")
    
    def _get_pdf_page_count(self, file_path: str) -> int:
        """Get the number of pages in a PDF."""
        images = pdf2image.convert_from_path(file_path)
        return len(images)
    
    def _get_docx_page_count(self, file_path: str) -> int:
        """Estimate the number of pages in a DOCX (simplified approach)."""
        doc = Document(file_path)
        # This is a rough estimation - DOCX doesn't have explicit page breaks
        # We'll count page breaks and estimate based on content
        page_breaks = 0
        for paragraph in doc.paragraphs:
            if paragraph.runs:
                for run in paragraph.runs:
                    if '\f' in run.text or '\x0c' in run.text:  # Page break characters
                        page_breaks += 1
        
        # If no explicit page breaks found, estimate based on content length
        if page_breaks == 0:
            total_paragraphs = len(doc.paragraphs)
            estimated_pages = max(1, (total_paragraphs + 19) // 20)  # Roughly 20 paragraphs per page
            return min(estimated_pages, 5)  # Cap at 5 for our sample
        
        return page_breaks + 1
    
    def _get_pptx_slide_count(self, file_path: str) -> int:
        """Get the number of slides in a PPTX."""
        prs = Presentation(file_path)
        return len(prs.slides)
    
    def _get_excel_sheet_count(self, file_path: str) -> int:
        """Get the number of sheets in an Excel file."""
        wb = openpyxl.load_workbook(file_path)
        return len(wb.worksheets)
    
    def _preview_pdf_page(self, file_path: str, page_number: int) -> Image.Image:
        """Generate preview for a PDF page."""
        images = pdf2image.convert_from_path(file_path, first_page=page_number, last_page=page_number, dpi=150)
        if images:
            return images[0]
        return self._create_error_image(f"PDF page {page_number} not found")
    
    def _preview_docx_page(self, file_path: str, page_number: int) -> Image.Image:
        """Generate preview for a DOCX page (simplified text rendering)."""
        doc = Document(file_path)
        
        # Create a white background image
        img_width, img_height = 800, 1000
        img = Image.new('RGB', (img_width, img_height), 'white')
        draw = ImageDraw.Draw(img)
        
        try:
            # Try to use a default font, fallback to default if not available
            font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
            font_text = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
        except:
            font_title = ImageFont.load_default()
            font_text = ImageFont.load_default()
        
        y_position = 50
        
        # Calculate which paragraphs belong to this page
        paragraphs_per_page = 20
        start_para = (page_number - 1) * paragraphs_per_page
        end_para = start_para + paragraphs_per_page
        
        # Add page header
        draw.text((50, 20), f"DOCX Document - Page {page_number}", fill='black', font=font_title)
        
        # Render paragraphs for this page
        for i, paragraph in enumerate(doc.paragraphs[start_para:end_para], start=start_para):
            if y_position > img_height - 100:
                break
            
            text = paragraph.text.strip()
            if text:
                # Wrap text to fit width
                wrapped_text = self._wrap_text(text, font_text, img_width - 100)
                for line in wrapped_text:
                    if y_position > img_height - 50:
                        break
                    draw.text((50, y_position), line, fill='black', font=font_text)
                    y_position += 25
                y_position += 10  # Extra space between paragraphs
        
        return img
    
    def _preview_pptx_slide(self, file_path: str, page_number: int) -> Image.Image:
        """Generate preview for a PPTX slide."""
        prs = Presentation(file_path)
        
        if page_number > len(prs.slides):
            return self._create_error_image(f"Slide {page_number} not found")
        
        slide = prs.slides[page_number - 1]
        
        # Create a white background image
        img_width, img_height = 800, 600
        img = Image.new('RGB', (img_width, img_height), 'white')
        draw = ImageDraw.Draw(img)
        
        try:
            font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 32)
            font_text = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
        except:
            font_title = ImageFont.load_default()
            font_text = ImageFont.load_default()
        
        y_position = 50
        
        # Extract and render slide content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Determine if this is likely a title (first text or larger)
                is_title = y_position == 50 or len(text) < 100
                current_font = font_title if is_title else font_text
                
                # Wrap and draw text
                wrapped_text = self._wrap_text(text, current_font, img_width - 100)
                for line in wrapped_text:
                    if y_position > img_height - 50:
                        break
                    draw.text((50, y_position), line, fill='black', font=current_font)
                    y_position += 40 if is_title else 25
                
                y_position += 20  # Extra space between text blocks
        
        return img
    
    def _preview_excel_sheet(self, file_path: str, page_number: int) -> Image.Image:
        """Generate preview for an Excel sheet."""
        wb = openpyxl.load_workbook(file_path)
        
        if page_number > len(wb.worksheets):
            return self._create_error_image(f"Sheet {page_number} not found")
        
        ws = wb.worksheets[page_number - 1]
        
        # Create a white background image
        img_width, img_height = 1000, 800
        img = Image.new('RGB', (img_width, img_height), 'white')
        draw = ImageDraw.Draw(img)
        
        try:
            font_header = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16)
            font_cell = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
        except:
            font_header = ImageFont.load_default()
            font_cell = ImageFont.load_default()
        
        # Draw sheet name
        draw.text((20, 20), f"Excel Sheet: {ws.title}", fill='black', font=font_header)
        
        # Draw grid and content
        start_x, start_y = 20, 60
        cell_width, cell_height = 150, 30
        max_rows, max_cols = 20, 6
        
        # Draw headers and grid
        for row in range(max_rows):
            for col in range(max_cols):
                x = start_x + col * cell_width
                y = start_y + row * cell_height
                
                # Draw cell border
                draw.rectangle([x, y, x + cell_width, y + cell_height], outline='black', width=1)
                
                # Get cell value
                excel_row = row + 1
                excel_col = col + 1
                cell = ws.cell(row=excel_row, column=excel_col)
                cell_value = str(cell.value) if cell.value is not None else ""
                
                # Truncate long text
                if len(cell_value) > 15:
                    cell_value = cell_value[:12] + "..."
                
                # Draw cell content
                text_x = x + 5
                text_y = y + 8
                current_font = font_header if row == 0 else font_cell
                draw.text((text_x, text_y), cell_value, fill='black', font=current_font)
        
        return img
    
    def _wrap_text(self, text: str, font, max_width: int) -> List[str]:
        """Wrap text to fit within specified width."""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = font.getbbox(test_line)
            text_width = bbox[2] - bbox[0]
            
            if text_width <= max_width:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    # Word is too long, break it
                    lines.append(word)
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return lines
    
    def _create_error_image(self, error_message: str) -> Image.Image:
        """Create an error image with the specified message."""
        img = Image.new('RGB', (800, 600), 'white')
        draw = ImageDraw.Draw(img)
        
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except:
            font = ImageFont.load_default()
        
        # Center the error message
        bbox = font.getbbox(error_message)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (800 - text_width) // 2
        y = (600 - text_height) // 2
        
        draw.text((x, y), error_message, fill='red', font=font)
        return img

